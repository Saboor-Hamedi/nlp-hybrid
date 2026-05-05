import io
import os

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches


def remove_watermarks_from_images(file_path, output_path):
    """Remove watermarks from images embedded in PowerPoint presentations."""
    prs = Presentation(file_path)

    # Get all image parts from the presentation
    total_removed = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"Processing slide {slide_num}...")

        for shape_num, shape in enumerate(slide.shapes):
            # Check if shape is a picture
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    image_bytes = image.blob

                    # Open image with PIL
                    img = Image.open(io.BytesIO(image_bytes))
                    print(f"  Found image: {img.size}, mode: {img.mode}")

                    # Convert to RGB if needed
                    if img.mode in ('RGBA', 'LA', 'P'):
                        img = img.convert('RGB')

                    # Remove watermark from bottom-right corner (20% of image)
                    width, height = img.size
                    watermark_x = int(width * 0.8)
                    watermark_y = int(height * 0.8)

                    # Get the corner region to analyze
                    corner_region = img.crop((watermark_x, watermark_y, width, height))

                    # Check if corner has light text (watermark)
                    corner_pixels = list(corner_region.getdata())
                    avg_brightness = sum(sum(p) if isinstance(p, tuple) else p for p in corner_pixels) / len(corner_pixels) / 3

                    print(f"    Corner brightness: {avg_brightness:.1f}")

                    # If corner has watermark (light area), inpaint it with surrounding color
                    if avg_brightness > 150:  # Light watermark detected
                        print(f"    Removing watermark from bottom-right...")

                        # Get background color from nearby area
                        bg_sample = img.crop((watermark_x - 50, watermark_y - 50, watermark_x, watermark_y))
                        bg_pixels = list(bg_sample.getdata())
                        avg_color = tuple(int(sum(p[i] for p in bg_pixels) / len(bg_pixels)) for i in range(3))

                        # Fill watermark area with background color
                        draw = ImageDraw.Draw(img)
                        draw.rectangle([(watermark_x, watermark_y), (width, height)], fill=avg_color)

                        # Save modified image back
                        img_byte_arr = io.BytesIO()
                        img.save(img_byte_arr, format='PNG')
                        img_byte_arr.seek(0)

                        # Replace image in presentation
                        image.blob = img_byte_arr.getvalue()
                        total_removed += 1
                        print(f"    ✓ Watermark removed from image")

                except Exception as e:
                    print(f"  Error processing image: {e}")

    prs.save(output_path)
    print(f"\n✓ Complete. Processed {total_removed} images. Saved as: {output_path}")


# Run the function
remove_watermarks_from_images("input_from_notebooklm.pptx", "no_watermark.pptx")
